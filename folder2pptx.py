"""
Build a PowerPoint (.pptx) slideshow from images and videos in a folder.
One slide per file, centered and scaled to fit the slide. Videos are timed to start when
the slide is shown (not click-to-play). Run: python folder2pptx.py -h
"""

from __future__ import annotations

import argparse
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

IMAGE_EXT = frozenset({".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".tif", ".tiff"})
VIDEO_EXT = frozenset({".mp4", ".mov", ".avi", ".webm", ".mkv", ".m4v", ".wmv"})
VIDEO_MIME = {
    ".mp4": "video/mp4",
    ".mov": "video/quicktime",
    ".avi": "video/x-msvideo",
    ".webm": "video/webm",
    ".mkv": "video/x-matroska",
    ".m4v": "video/mp4",
    ".wmv": "video/x-ms-wmv",
}


def _natural_key(s: str) -> list:
    """Sort strings with embedded numbers numerically (file2 before file10)."""
    parts = re.split(r"(\d+)", s.lower())
    return [int(p) if p.isdigit() else p for p in parts]


def collect_media(folder: Path, recursive: bool) -> list[Path]:
    pattern = "**/*" if recursive else "*"
    files: list[Path] = []
    for p in folder.glob(pattern):
        if not p.is_file():
            continue
        ext = p.suffix.lower()
        if ext in IMAGE_EXT or ext in VIDEO_EXT:
            files.append(p)
    return files


def _blank_slide_layout(prs):
    for layout in prs.slide_layouts:
        if (layout.name or "").lower() == "blank":
            return layout
    if len(prs.slide_layouts) > 6:
        return prs.slide_layouts[6]
    return prs.slide_layouts[-1]


def _fit_picture_on_slide(slide, prs, image_path: Path) -> None:
    slide_w, slide_h = prs.slide_width, prs.slide_height
    pic = slide.shapes.add_picture(str(image_path), 0, 0)
    iw, ih = pic.width, pic.height
    if iw <= 0 or ih <= 0:
        return
    scale = min(slide_w / iw, slide_h / ih)
    pic.width = int(iw * scale)
    pic.height = int(ih * scale)
    pic.left = (slide_w - pic.width) // 2
    pic.top = (slide_h - pic.height) // 2


def _ffprobe_video_size(path: Path) -> tuple[int, int] | None:
    exe = shutil.which("ffprobe")
    if not exe:
        return None
    try:
        r = subprocess.run(
            [
                exe,
                "-v",
                "error",
                "-select_streams",
                "v:0",
                "-show_entries",
                "stream=width,height",
                "-of",
                "csv=p=0:s=x",
                str(path),
            ],
            capture_output=True,
            text=True,
            timeout=120,
        )
    except (OSError, subprocess.TimeoutExpired):
        return None
    if r.returncode != 0 or not (r.stdout or "").strip():
        return None
    try:
        w, h = r.stdout.strip().split("x", 1)
        wi, hi = int(w), int(h)
        if wi > 0 and hi > 0:
            return wi, hi
    except ValueError:
        pass
    return None


def _ffmpeg_first_frame(video_path: Path, out_path: Path) -> bool:
    exe = shutil.which("ffmpeg")
    if not exe:
        return False
    try:
        r = subprocess.run(
            [
                exe,
                "-y",
                "-ss",
                "0",
                "-i",
                str(video_path),
                "-frames:v",
                "1",
                "-q:v",
                "2",
                str(out_path),
            ],
            capture_output=True,
            timeout=120,
        )
    except (OSError, subprocess.TimeoutExpired):
        return False
    return r.returncode == 0 and out_path.is_file()


def _video_dimensions(path: Path) -> tuple[int, int]:
    try:
        import cv2

        if not hasattr(cv2, "VideoCapture"):
            raise AttributeError("cv2.VideoCapture missing (broken opencv install?)")
        cap = cv2.VideoCapture(str(path))
        try:
            if not cap.isOpened():
                raise OSError("VideoCapture did not open file")
            w = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH) or 0)
            h = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT) or 0)
            if w > 0 and h > 0:
                return w, h
        finally:
            cap.release()
    except Exception:
        pass
    wh = _ffprobe_video_size(path)
    if wh:
        return wh
    return 1920, 1080


def _extract_video_poster(video_path: Path, out_path: Path) -> bool:
    try:
        import cv2

        if not hasattr(cv2, "VideoCapture"):
            raise AttributeError("cv2.VideoCapture missing")
        cap = cv2.VideoCapture(str(video_path))
        try:
            ok, frame = cap.read()
            if ok and frame is not None:
                return bool(cv2.imwrite(str(out_path), frame))
        finally:
            cap.release()
    except Exception:
        pass
    return _ffmpeg_first_frame(video_path, out_path)


def _wrap_root_child_videos_in_seq(sld) -> None:
    """PowerPoint often expects p:video under p:seq; python-pptx puts video directly in childTnLst."""
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls, qn

    ctl_list = sld.xpath("./p:timing/p:tnLst/p:par/p:cTn/p:childTnLst")
    if not ctl_list:
        return
    ctl = ctl_list[0]
    videos = [el for el in list(ctl) if el.tag == qn("p:video")]
    if not videos:
        return
    ids = [int(x) for x in sld.xpath(".//p:cTn/@id")]
    next_id = (max(ids) if ids else 0) + 1
    for video in videos:
        ctl.remove(video)
        seq = parse_xml(
            (
                "<p:seq %s concurrent=\"1\" nextAc=\"seek\">\n"
                '  <p:cTn id="%d" dur="indefinite" restart="never" nodeType="seq">\n'
                "    <p:childTnLst/>\n"
                "  </p:cTn>\n"
                "</p:seq>\n"
            )
            % (nsdecls("p"), next_id)
        )
        next_id += 1
        # p:seq from parse_xml is not a registered Oxml type — xpath("p:...") lacks ns bindings.
        ctn_el = seq.find(qn("p:cTn"))
        inner = ctn_el.find(qn("p:childTnLst")) if ctn_el is not None else None
        if inner is None:
            raise RuntimeError("internal: malformed p:seq wrapper")
        inner.append(video)
        ctl.append(seq)


def _patch_all_slide_videos_autoplay(sld) -> None:
    """Replace click-to-play (delay=indefinite, display=0) with start-on-slide (delay=0)."""
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls, qn

    for video in sld.xpath(".//p:timing//p:video"):
        cmedia = video.find(qn("p:cMediaNode"))
        if cmedia is None:
            continue
        ctn = cmedia.find(qn("p:cTn"))
        if ctn is None:
            continue
        ctn.attrib.pop("display", None)
        stc = ctn.find(qn("p:stCondLst"))
        if stc is not None:
            for child in list(stc):
                stc.remove(child)
        else:
            stc = parse_xml("<p:stCondLst %s/>" % nsdecls("p"))
            ctn.insert(0, stc)
        stc.append(parse_xml("<p:cond %s delay=\"0\"/>" % nsdecls("p")))


def _configure_video_slide_timing(slide) -> None:
    _wrap_root_child_videos_in_seq(slide._element)
    _patch_all_slide_videos_autoplay(slide._element)


def _add_video_slide(slide, prs, video_path: Path, tmpdir: Path) -> None:
    slide_w, slide_h = prs.slide_width, prs.slide_height
    vw, vh = _video_dimensions(video_path)
    scale = min(slide_w / vw, slide_h / vh)
    w = int(vw * scale)
    h = int(vh * scale)
    left = (slide_w - w) // 2
    top = (slide_h - h) // 2
    mime = VIDEO_MIME.get(video_path.suffix.lower(), "video/mp4")
    poster = tmpdir / f"poster_{video_path.stem}.jpg"
    poster_arg = str(poster) if _extract_video_poster(video_path, poster) else None
    slide.shapes.add_movie(
        str(video_path),
        left,
        top,
        w,
        h,
        poster_frame_image=poster_arg,
        mime_type=mime,
    )
    _configure_video_slide_timing(slide)


def build_pptx(
    folder: Path,
    output: Path,
    *,
    recursive: bool = False,
    sort_by: str = "name",
    verbose: bool = False,
) -> int:
    try:
        from pptx import Presentation
    except ImportError:
        print("Error: python-pptx is required. Install with: pip install python-pptx", file=sys.stderr)
        return 1

    if not folder.is_dir():
        print(f"Error: Not a directory: {folder}", file=sys.stderr)
        return 1

    files = collect_media(folder, recursive)
    if not files:
        print(f"No images or videos found in {folder}", file=sys.stderr)
        return 1

    if sort_by == "mtime":
        files.sort(key=lambda p: p.stat().st_mtime)
    else:
        files.sort(key=lambda p: _natural_key(p.name))

    prs = Presentation()
    layout = _blank_slide_layout(prs)

    with tempfile.TemporaryDirectory(prefix="folder2pptx_") as tmp:
        tmpdir = Path(tmp)
        for path in files:
            ext = path.suffix.lower()
            try:
                if ext in IMAGE_EXT:
                    slide = prs.slides.add_slide(layout)
                    _fit_picture_on_slide(slide, prs, path)
                elif ext in VIDEO_EXT:
                    slide = prs.slides.add_slide(layout)
                    _add_video_slide(slide, prs, path, tmpdir)
                else:
                    continue
                if verbose:
                    print(path)
            except Exception as e:
                print(f"Warning: skipped {path}: {e}", file=sys.stderr)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    print(f"Wrote {len(prs.slides)} slide(s) to {output}")
    if any(p.suffix.lower() in VIDEO_EXT for p in files):
        print(
            "Tip: Test embedded video with Slide Show (F5). PowerPoint’s codecs differ from VLC/Edge; "
            "re-encode only if slideshow playback still fails.",
            file=sys.stderr,
        )
    return 0


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Create a .pptx slideshow from images and videos in a folder (one slide per file).",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python folder2pptx.py -d ./photos -o album.pptx
  python folder2pptx.py -d ./media -r -o show.pptx
  python folder2pptx.py -d ./clips --sort mtime -o timeline.pptx

Supported images: jpg, png, gif, webp, bmp, tiff
Supported videos: mp4, mov, avi, webm, mkv, m4v, wmv
        """,
    )
    parser.add_argument("-d", "--directory", required=True, help="Folder containing images and/or videos")
    parser.add_argument("-o", "--output", default="slideshow.pptx", help="Output .pptx path (default: slideshow.pptx)")
    parser.add_argument("-r", "--recursive", action="store_true", help="Include subfolders")
    parser.add_argument(
        "--sort",
        dest="sort_by",
        choices=("name", "mtime"),
        default="name",
        help="Order slides by filename (natural sort) or file modification time (default: name)",
    )
    parser.add_argument("-v", "--verbose", action="store_true", help="Print each file as it is added")
    args = parser.parse_args()
    sys.exit(build_pptx(Path(args.directory), Path(args.output), recursive=args.recursive, sort_by=args.sort_by, verbose=args.verbose))


if __name__ == "__main__":
    main()
